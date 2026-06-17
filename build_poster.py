"""Botany 2026 poster for Flora of the World — 48" x 36".

Zones (top → bottom):
  HEADER          title, authors, FotW logo, Heliconia banner accent, URL pill
  STAT STRIP      seven headline numbers
  NARRATIVE       3 columns: Mission · Approach (with clear 12 callout) · Community
  FOUNDER LEGACY  3 panels: Davidson exploration · Library archive · Living gardens
  GEOGRAPHIC      full-width occurrence map
  DATA            3 panels: top-12 countries · top-15 families · EDGE plants
  PARTNERS + CTA  WFO + EDGE partnership · BSU Endowed Chair · JOIN US + QR
  FOOTER          References (incl. publications URL) + Contact
"""

import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = "/Users/sven/Documents/Current_projects/Endowed_Chair/Botany_2026"
FIG_DIR = os.path.join(HERE, "figures")
def F(name): return os.path.join(FIG_DIR, name)

# ------------------------------------------------------------- palette
# FotW brand-inspired blue palette (the site uses #2297F2 in its navbar)
GREEN_DARK = RGBColor(0x0B, 0x39, 0x54)   # deep navy — headers/dark bands
GREEN_MID  = RGBColor(0x15, 0x68, 0xB0)   # mid blue — dividers/secondary
SAGE       = RGBColor(0x75, 0xCF, 0xF0)   # sky blue — soft accent stripe
GOLD       = RGBColor(0xF5, 0x8E, 0x1F)   # amber — warm pop / underline
CREAM      = RGBColor(0xF4, 0xF8, 0xFB)   # cool off-white background
CARD       = RGBColor(0xFF, 0xFF, 0xFF)
INK        = RGBColor(0x1A, 0x2A, 0x3A)   # cool dark text
INK_SOFT   = RGBColor(0x5A, 0x68, 0x78)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GOLD_SOFT  = RGBColor(0xF8, 0xDD, 0xB3)   # warm sand — light accents on dark
BLUE_BRAND = RGBColor(0x22, 0x97, 0xF2)   # FotW navbar blue — bright accent

SERIF = "Cambria"
SANS  = "Calibri"

prs = Presentation()
prs.slide_width  = Inches(48)
prs.slide_height = Inches(36)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# ------------------------------------------------------------- helpers
from pptx.oxml.ns import qn

def _strip_empty_txbody(shp):
    """Remove auto-generated empty txBody from a shape — PowerPoint flags
    rectangles whose auto txBody contains a paragraph with no run."""
    tb = shp._element.find(qn("p:txBody"))
    if tb is not None:
        shp._element.remove(tb)

def add_rect(left, top, width, height, fill, line=None):
    if width <= 0:  width  = 0.10
    if height <= 0: height = 0.10
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    _strip_empty_txbody(shp)
    return shp

def add_text(left, top, width, height, text, *, font=SERIF, size=20,
             bold=False, italic=False, color=INK, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.15, letter_spacing=None):
    # Defensive: any layout bug that computes a negative or near-zero height
    # produces a malformed shape; clamp to a small positive value.
    if width <= 0:  width  = 0.10
    if height <= 0: height = 0.10
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run(); r.text = line
        f = r.font; f.name = font; f.size = Pt(size); f.bold = bold
        f.italic = italic; f.color.rgb = color
        if letter_spacing is not None:
            rPr = r._r.get_or_add_rPr()
            rPr.set('spc', str(int(letter_spacing * 100)))
    return tb

def add_paragraphs(left, top, width, height, blocks, *, font=SERIF, size=20,
                   color=INK, align=PP_ALIGN.LEFT, line_spacing=1.2,
                   space_after=8):
    if width <= 0:  width  = 0.10
    if height <= 0: height = 0.10
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    for i, item in enumerate(blocks):
        text, opts = item if isinstance(item, tuple) else (item, {})
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get("align", align)
        p.line_spacing = opts.get("line_spacing", line_spacing)
        p.space_after = Pt(opts.get("space_after", space_after))
        runs = opts.get("runs") or [(text, {})]
        for rt, ro in runs:
            r = p.add_run(); r.text = rt
            f = r.font
            f.name = ro.get("font", opts.get("font", font))
            f.size = Pt(ro.get("size", opts.get("size", size)))
            f.bold = ro.get("bold", opts.get("bold", False))
            f.italic = ro.get("italic", opts.get("italic", False))
            f.color.rgb = ro.get("color", opts.get("color", color))
    return tb

def add_line(x1, y1, x2, y2, color=GREEN_MID, weight=1.5):
    """Draw a line as a thin filled rectangle (more compatible than connector)."""
    # Convert pt weight to inches (1 pt = 1/72 in)
    thickness = max(weight / 72.0, 0.012)
    dx = x2 - x1; dy = y2 - y1
    if abs(dx) >= abs(dy):
        # Horizontal-leaning
        left = min(x1, x2); width = abs(dx) if dx else thickness
        top = min(y1, y2) - thickness / 2; height = thickness
    else:
        left = min(x1, x2) - thickness / 2; width = thickness
        top = min(y1, y2); height = abs(dy) if dy else thickness
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    _strip_empty_txbody(shp)
    return shp

def section_header(left, top, width, label, color=GREEN_DARK, accent=GOLD, size=24):
    add_text(left, top, width, 0.55, label.upper(), font=SANS, size=size,
             bold=True, color=color, letter_spacing=4)
    add_line(left, top + 0.65, left + 1.4, top + 0.65, color=accent, weight=3)

def add_image_fit(path, left, top, width, height, *, frame_color=None,
                  frame_weight=1.0, anchor="center"):
    """Place image scaled to FILL the box (crop to aspect)."""
    im = Image.open(path)
    iw, ih = im.size
    box_ratio = width / height
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        new_iw = int(ih * box_ratio)
        if anchor == "left": x0 = 0
        elif anchor == "right": x0 = iw - new_iw
        else: x0 = (iw - new_iw) // 2
        cropped = im.crop((x0, 0, x0 + new_iw, ih))
    else:
        new_ih = int(iw / box_ratio)
        if anchor == "top": y0 = 0
        elif anchor == "bottom": y0 = ih - new_ih
        else: y0 = (ih - new_ih) // 2
        cropped = im.crop((0, y0, iw, y0 + new_ih))
    tmp = path.replace(".jpg", "_fit.jpg").replace(".jpeg", "_fit.jpg").replace(".png", "_fit.png")
    cropped.save(tmp, quality=92)
    pic = slide.shapes.add_picture(tmp, Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    if frame_color is not None:
        ln = pic.line; ln.color.rgb = frame_color; ln.width = Pt(frame_weight)
    return pic

def add_image_contain(path, left, top, width, height):
    """Place image preserving aspect, centered inside the box."""
    im = Image.open(path)
    iw, ih = im.size
    box_ratio = width / height
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        new_w = width; new_h = width / img_ratio
        x = left; y = top + (height - new_h) / 2
    else:
        new_h = height; new_w = height * img_ratio
        x = left + (width - new_w) / 2; y = top
    return slide.shapes.add_picture(path, Inches(x), Inches(y),
                                    Inches(new_w), Inches(new_h))

# ============================================================ background
add_rect(0, 0, 48, 36, CREAM)

# ============================================================ HEADER
HDR_H = 4.4
add_rect(0, 0, 48, HDR_H, GREEN_DARK)
add_rect(0, HDR_H, 48, 0.10, GOLD)
add_rect(0, 0, 0.35, HDR_H, SAGE)

# FotW logo (white) top-left
fotw_h = 1.8
fotw_w = fotw_h * (1600 / 957)
slide.shapes.add_picture(F("fotw_logo.png"), Inches(0.9), Inches(0.40),
                         Inches(fotw_w), Inches(fotw_h))

# Overline beside logo
add_text(0.9 + fotw_w + 0.5, 0.50, 18, 0.45,
         "BOTANY 2026  ·  POSTER PRESENTATION",
         font=SANS, size=17, color=GOLD_SOFT, letter_spacing=10, bold=True)

# Title — even narrower so the three banner photos can each be big
title_left = 0.9 + fotw_w + 0.5
title_w    = 20
add_text(title_left, 0.95, title_w, 1.5, "Flora of the World",
         font=SERIF, size=50, bold=True, color=WHITE, line_spacing=1.0)
add_text(title_left, 2.20, title_w, 0.60,
         "An online platform documenting plants in their native habitats",
         font=SERIF, size=18, italic=True, color=GOLD_SOFT, line_spacing=1.05)

# Authors + affiliations
add_text(title_left, 2.90, title_w, 0.45,
         "Sven Buerki¹²  ·  Kenny Williams¹  ·  Jim Beck²  ·  Sharon R. Christoph¹",
         font=SANS, size=15, color=WHITE, line_spacing=1.0)
add_text(title_left, 3.40, title_w + 5, 0.45,
         "¹ Flora of the World Foundation     ² Boise State University, Idaho, USA",
         font=SANS, size=11, italic=True, color=GOLD_SOFT)

# Banner accent: three species photos side-by-side, filling header height
# Freycinetia biloba (4:3 landscape) · Piper auritum (3:4 portrait) ·
# Carnegiea gigantea (4:3 landscape)
banner_top = 0.35
banner_h   = 3.70                 # fills most of the 4.4" header
gap        = 0.15
# Photo widths, keeping each source aspect (no crop)
frey_w  = banner_h * 4 / 3        # 4:3 landscape
piper_w = banner_h * 3 / 4        # 3:4 portrait
carn_w  = banner_h * 4 / 3        # 4:3 landscape
# Layout: title ends ~24.4 (title_left + title_w), URL pill starts at 38.6
banner_left = 25.4                # gap to title and to URL pill

def banner_photo(left, top, w, h, image_path, runs):
    add_rect(left - 0.07, top - 0.07, w + 0.14, h + 0.14, GOLD)
    add_image_fit(F(image_path), left, top, w, h, anchor="center")
    add_rect(left, top + h - 0.28, w, 0.28, GREEN_DARK)
    add_paragraphs(
        left, top + h - 0.26, w, 0.24,
        [("", {"runs": runs, "size": 8, "color": WHITE,
                "align": PP_ALIGN.CENTER, "space_after": 0})],
        font=SANS, line_spacing=1.0,
    )

# Freycinetia
banner_photo(banner_left, banner_top, frey_w, banner_h,
             "banner_flower.jpg",
             [("Freycinetia biloba", {"italic": True}),
              (" B.C.Stone  ·  Pandanaceae", {})])
# Piper
piper_left = banner_left + frey_w + gap
banner_photo(piper_left, banner_top, piper_w, banner_h,
             "banner_piper.jpg",
             [("Piper auritum", {"italic": True}),
              (" Kunth  ·  Piperaceae", {})])
# Carnegiea
carn_left = piper_left + piper_w + gap
banner_photo(carn_left, banner_top, carn_w, banner_h,
             "banner_carnegiea.jpg",
             [("Carnegiea gigantea", {"italic": True}),
              (" (Engelm.) Britton & Rose  ·  Cactaceae", {})])

# Right-side: URL pill with QR code on the right inside the same card
pill_left = 38.6
pill_top  = 0.50
pill_w    = 8.2
pill_h    = 3.40
add_rect(pill_left, pill_top, pill_w, pill_h, WHITE)
# Left part of pill: URL info
url_w = pill_w * 0.62
add_text(pill_left + 0.15, pill_top + 0.30, url_w - 0.30, 0.40,
         "VISIT THE PLATFORM",
         font=SANS, size=12, color=GREEN_MID,
         align=PP_ALIGN.CENTER, letter_spacing=6, bold=True)
add_text(pill_left + 0.15, pill_top + 0.75, url_w - 0.30, 0.70,
         "floraoftheworld.org",
         font=SERIF, size=24, color=GREEN_DARK, bold=True,
         align=PP_ALIGN.CENTER)
add_text(pill_left + 0.15, pill_top + 1.60, url_w - 0.30, 0.50,
         "Scan the QR code\nto explore the platform",
         font=SANS, size=13, italic=True, color=GREEN_MID,
         align=PP_ALIGN.CENTER, line_spacing=1.2)
add_text(pill_left + 0.15, pill_top + 2.55, url_w - 0.30, 0.65,
         "232,435 images · 6,010 species\n52 countries · 956 contributors",
         font=SANS, size=12, italic=True, color=INK_SOFT,
         align=PP_ALIGN.CENTER, line_spacing=1.25)
# Right part of pill: QR code square
qr_box_size = pill_h - 0.40
qr_box_left = pill_left + url_w + 0.10
slide.shapes.add_picture(F("qr_fotw.png"),
                         Inches(qr_box_left), Inches(pill_top + 0.20),
                         Inches(qr_box_size), Inches(qr_box_size))

# ============================================================ STAT STRIP
SS_TOP = HDR_H + 0.30
SS_H   = 1.55
SS_MARGIN = 0.9
ss_left = SS_MARGIN
ss_right = 48 - SS_MARGIN
stats_strip = [
    ("232,435", "DIGITAL IMAGES"),
    ("19,640",  "OCCURRENCES"),
    ("6,010+",  "SPECIES"),
    ("4,516",   "GENERA"),
    ("482",     "FAMILIES"),
    ("52",      "COUNTRIES"),
    ("46",      "HERBARIA"),
]
seg_w = (ss_right - ss_left) / len(stats_strip)
add_rect(ss_left, SS_TOP, ss_right - ss_left, SS_H, CARD)
for i, (big, lab) in enumerate(stats_strip):
    sx = ss_left + i * seg_w
    if i > 0:
        add_line(sx, SS_TOP + 0.20, sx, SS_TOP + SS_H - 0.20, color=SAGE, weight=1.0)
    add_text(sx, SS_TOP + 0.20, seg_w, 0.80, big, font=SERIF, size=42,
             bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER, line_spacing=1.0)
    add_text(sx, SS_TOP + 1.00, seg_w, 0.45, lab, font=SANS, size=14,
             bold=True, color=INK_SOFT, align=PP_ALIGN.CENTER,
             letter_spacing=4, line_spacing=1.0)
add_rect(ss_left, SS_TOP, ss_right - ss_left, 0.10, GOLD)

# ============================================================ NARRATIVE
NA_TOP = SS_TOP + SS_H + 0.35
NA_BOT = 9.80
NA_H = NA_BOT - NA_TOP

GAP = 0.55
COL_W = (48 - 2*SS_MARGIN - 2*GAP) / 3
COL1_L = SS_MARGIN
COL2_L = SS_MARGIN + COL_W + GAP
COL3_L = SS_MARGIN + 2*(COL_W + GAP)

# --- Col 1: MISSION
y = NA_TOP
section_header(COL1_L, y, COL_W, "Our Mission", size=22)
y += 0.80
add_text(COL1_L, y, COL_W, NA_H - 0.85,
         "Flora of the World documents the diversity of flowering "
         "plant families in their native environments, with a special "
         "focus on global biodiversity hotspots and botanical garden "
         "inventories — pairing high-resolution digital imagery with "
         "field-collected specimens to preserve morphological, ecological, "
         "and geographic context.",
         font=SERIF, size=19, color=INK, line_spacing=1.32)

# --- Col 2: INNOVATIVE APPROACH (with clearer 12 callout)
y = NA_TOP
section_header(COL2_L, y, COL_W, "An Innovative Approach", size=22)
y += 0.80
add_text(COL2_L, y, COL_W, 1.55,
         "What makes Flora of the World singular is the depth of its "
         "visual documentation: each occurrence is photographed many "
         "times to capture diagnostic features, habitat, and variation. "
         "We aim to document at least one species from every flowering "
         "plant family.",
         font=SERIF, size=17, color=INK, line_spacing=1.28)
y += 1.60
# Clear, labelled 12 callout — sized to fit the available cb_h
cb_h = max(NA_BOT - y - 0.05, 1.0)
add_rect(COL2_L, y, COL_W, cb_h, GREEN_DARK)
add_text(COL2_L + 0.30, y, 3.4, cb_h, "12",
         font=SERIF, size=96, bold=True, color=GOLD,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
         line_spacing=1.0)
# Right side labels — two-line stack, vertically centred in the callout
lx = COL2_L + 3.8
lw = COL_W - 4.0
label_h = cb_h * 0.40
gap = 0.05
total_lh = 2 * label_h + gap
top_offset = (cb_h - total_lh) / 2
add_text(lx, y + top_offset, lw, label_h,
         "PHOTOGRAPHS",
         font=SANS, size=22, bold=True, color=WHITE,
         letter_spacing=6, line_spacing=1.0,
         anchor=MSO_ANCHOR.MIDDLE)
add_text(lx, y + top_offset + label_h + gap, lw, label_h,
         "PER OCCURRENCE",
         font=SANS, size=18, bold=True, color=GOLD_SOFT,
         letter_spacing=6, line_spacing=1.0,
         anchor=MSO_ANCHOR.MIDDLE)

# --- Col 3: COMMUNITY — text + stat strip sized to fit available space
y = NA_TOP
section_header(COL3_L, y, COL_W, "A Global Community", size=22)
y += 0.80
add_text(COL3_L, y, COL_W, 1.05,
         "Flora of the World is built by people — collector profiles, "
         "a publications database, and pages for partner herbaria, gardens, "
         "and research organisations.",
         font=SERIF, size=16, color=INK, line_spacing=1.26)
y += 1.10
strip_h = max(NA_BOT - y - 0.05, 0.10)
seg_w3 = (COL_W - 0.30) / 3
items = [("956", "contributors"), ("91", "institutions"), ("138", "publications")]
# Render numbers + labels INSIDE the strip box, scaled to fit strip_h
num_h  = strip_h * 0.60
lab_h  = strip_h - num_h - 0.05
num_sz = max(18, int(num_h * 56))   # ~56pt per inch ≈ visually proportional
for i, (n, lab) in enumerate(items):
    sx = COL3_L + i * (seg_w3 + 0.15)
    add_rect(sx, y, seg_w3, strip_h, GREEN_MID)
    add_text(sx, y + 0.05, seg_w3, num_h, n,
             font=SERIF, size=num_sz, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, line_spacing=1.0,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(sx, y + 0.05 + num_h, seg_w3, lab_h, lab,
             font=SANS, size=12, italic=True, color=GOLD_SOFT,
             align=PP_ALIGN.CENTER, line_spacing=1.0,
             anchor=MSO_ANCHOR.MIDDLE)

# ============================================================ FOUNDER LEGACY
FL_TOP = NA_BOT + 0.35
FL_BOT = 14.35
FL_H = FL_BOT - FL_TOP

# Title strip across full band
add_text(SS_MARGIN, FL_TOP, 38, 0.55,
         "A LEGACY OF EXPLORATION — Dr. Christopher Davidson (1944–2022)",
         font=SANS, size=24, bold=True, color=GREEN_DARK, letter_spacing=4)
add_line(SS_MARGIN, FL_TOP + 0.65, SS_MARGIN + 1.4, FL_TOP + 0.65,
         color=GOLD, weight=3)

# Three card panels — compact layout for 4.2" band height
fl_inner_top = FL_TOP + 0.85
fl_inner_h   = FL_BOT - fl_inner_top
panel_w = (48 - 2*SS_MARGIN - 2*GAP) / 3
p1_l = SS_MARGIN
p2_l = SS_MARGIN + panel_w + GAP
p3_l = SS_MARGIN + 2*(panel_w + GAP)

# ---- PANEL 1: Davidson — full 3:4 portrait (contain, NO crop) on the left + text
add_rect(p1_l, fl_inner_top, panel_w, fl_inner_h, CARD)
# Photo box matches the source 3:4 aspect so the whole figure shows uncropped.
ph_pad = 0.18
ph_h = fl_inner_h - 2 * ph_pad
ph_w = ph_h * 0.75           # 3:4 portrait aspect — head to feet preserved
ph_x = p1_l + ph_pad
ph_y = fl_inner_top + ph_pad
add_image_contain(F("davidson_field.jpg"), ph_x, ph_y, ph_w, ph_h)
# Text on the right
tx_l = ph_x + ph_w + 0.30
tx_w = panel_w - (tx_l - p1_l) - 0.25
add_text(tx_l, fl_inner_top + 0.20, tx_w, 0.45,
         "EXPLORATION",
         font=SANS, size=15, bold=True, color=GOLD,
         letter_spacing=8)
add_paragraphs(
    tx_l, fl_inner_top + 0.75, tx_w, fl_inner_h - 1.50,
    [
        ("", {"runs": [
            ("Founded by ", {}),
            ("Dr. Christopher Davidson", {"bold": True}),
            (", an Idaho-born botanist who, with Sharon Christoph, "
             "travelled to ", {}),
            ("45 countries over 25 years", {"bold": True}),
            (" to document flowering plants in situ and to support "
             "botanical research and conservation worldwide ", {}),
            ("(Taylor et al. 2023)", {"italic": True, "color": GREEN_MID}),
            (".", {}),
        ], "size": 16, "color": INK, "space_after": 6}),
    ],
    font=SERIF, line_spacing=1.28,
)
add_text(tx_l, fl_inner_top + fl_inner_h - 0.65, tx_w, 0.55,
         "Dr. Davidson collecting Pancheria hirsuta (Cunoniaceae), "
         "New Caledonia.",
         font=SANS, size=12, italic=True, color=INK_SOFT, line_spacing=1.18)

# ---- PANEL 2: Library — pre-cropped 1.66:1 photo, fit-fill to fill the box
add_rect(p2_l, fl_inner_top, panel_w, fl_inner_h, CARD)
add_text(p2_l + 0.25, fl_inner_top + 0.20, panel_w - 0.50, 0.45,
         "THE ARCHIVE",
         font=SANS, size=15, bold=True, color=GOLD, letter_spacing=8)
lib_top = fl_inner_top + 0.75
lib_w_avail = panel_w - 0.40
lib_h = fl_inner_h - 0.75 - 1.00
add_image_fit(F("library.jpg"), p2_l + 0.20, lib_top, lib_w_avail, lib_h,
              frame_color=GREEN_DARK, frame_weight=1.0)
add_text(p2_l + 0.25, lib_top + lib_h + 0.10, panel_w - 0.50, 0.85,
         "Dr. Davidson assembled an extensive natural-history library — "
         "taxonomic and ecological context being integrated into the platform.",
         font=SERIF, size=14, italic=True, color=INK, line_spacing=1.22)

# ---- PANEL 3: Living Legacy — caption on LEFT, gardens (bigger) on RIGHT
add_rect(p3_l, fl_inner_top, panel_w, fl_inner_h, CARD)
add_text(p3_l + 0.25, fl_inner_top + 0.20, panel_w - 0.50, 0.45,
         "LIVING LEGACY",
         font=SANS, size=15, bold=True, color=GOLD, letter_spacing=8)
# Left column: short caption text
txt_l_local = p3_l + 0.25
txt_w_local = panel_w * 0.30
add_text(txt_l_local, fl_inner_top + 0.75, txt_w_local,
         fl_inner_h - 1.00,
         "Two private gardens (Idaho) founded by Dr. Davidson — living "
         "plant collections carrying the mission forward.",
         font=SERIF, size=14, italic=True, color=INK, line_spacing=1.22)
# Right area: two garden photos side-by-side, BIGGER (fit-fill so they fill)
gd_top = fl_inner_top + 0.75
gd_avail_h = fl_inner_h - 0.75 - 0.50   # leave a small caption strip below
photos_left_x = txt_l_local + txt_w_local + 0.30
photos_avail_w = (p3_l + panel_w - 0.25) - photos_left_x
gap_px = 0.18
each_w = (photos_avail_w - gap_px) / 2
each_h = gd_avail_h
add_image_fit(F("boise_garden.jpg"), photos_left_x, gd_top,
              each_w, each_h, frame_color=GREEN_DARK, frame_weight=1.0)
add_image_fit(F("mccall_garden.jpg"),
              photos_left_x + each_w + gap_px, gd_top,
              each_w, each_h, frame_color=GREEN_DARK, frame_weight=1.0)
# Labels under each photo
lbl_y = gd_top + each_h + 0.05
add_text(photos_left_x, lbl_y, each_w, 0.32, "Boise garden",
         font=SANS, size=12, italic=True, color=INK_SOFT,
         align=PP_ALIGN.CENTER)
add_text(photos_left_x + each_w + gap_px, lbl_y, each_w, 0.32,
         "McCall garden", font=SANS, size=12, italic=True,
         color=INK_SOFT, align=PP_ALIGN.CENTER)

# ============================================================ GEOGRAPHIC REACH
GEO_TOP = FL_BOT + 0.30
GEO_BOT = 25.50
GEO_H = GEO_BOT - GEO_TOP
geo_left = SS_MARGIN
geo_right = 48 - SS_MARGIN
geo_w = geo_right - geo_left

add_text(geo_left, GEO_TOP, 24, 0.55, "GEOGRAPHIC REACH",
         font=SANS, size=24, bold=True, color=GREEN_DARK, letter_spacing=4)
add_text(geo_left + 14, GEO_TOP + 0.05, geo_w - 14, 0.55,
         "19,640 occurrences  ·  52 countries  ·  722 locations  ·  3,842 collection events",
         font=SERIF, size=18, italic=True, color=INK_SOFT,
         align=PP_ALIGN.RIGHT)
add_line(geo_left, GEO_TOP + 0.65, geo_left + 1.4, GEO_TOP + 0.65,
         color=GOLD, weight=3)

map_top = GEO_TOP + 0.85
map_h = GEO_BOT - map_top - 0.05
add_image_fit(F("map_full.png"), geo_left, map_top, geo_w, map_h,
              frame_color=GREEN_DARK, frame_weight=1.5)

# ============================================================ DATA & DIVERSITY
DD_TOP = GEO_BOT + 0.25
DD_BOT = 31.60
DD_H = DD_BOT - DD_TOP

add_text(geo_left, DD_TOP, 30, 0.55, "DATA · DIVERSITY · CONSERVATION",
         font=SANS, size=24, bold=True, color=GREEN_DARK, letter_spacing=4)
add_line(geo_left, DD_TOP + 0.65, geo_left + 1.4, DD_TOP + 0.65,
         color=GOLD, weight=3)

dd_top = DD_TOP + 0.90
dd_h = DD_BOT - dd_top
sub_w = (geo_w - 2 * GAP) / 3
# Inner panel layout: small title + small subtitle, then big chart image
HEAD_H = 0.85  # title + subtitle area inside each panel

# Panel 1: countries (left)
add_rect(geo_left, dd_top, sub_w, dd_h, CARD)
add_text(geo_left + 0.20, dd_top + 0.12, sub_w - 0.40, 0.42,
         "TOP 12 COUNTRIES", font=SANS, size=18, bold=True,
         color=GREEN_DARK, letter_spacing=4)
add_text(geo_left + 0.20, dd_top + 0.50, sub_w - 0.40, 0.38,
         "Where the collection has been built",
         font=SERIF, size=14, italic=True, color=INK_SOFT)
add_image_contain(F("chart_countries.png"),
                  geo_left + 0.15, dd_top + HEAD_H,
                  sub_w - 0.30, dd_h - HEAD_H - 0.20)

# Panel 2: families (middle)
p2 = geo_left + sub_w + GAP
add_rect(p2, dd_top, sub_w, dd_h, CARD)
add_text(p2 + 0.20, dd_top + 0.12, sub_w - 0.40, 0.42,
         "TOP 15 FAMILIES", font=SANS, size=18, bold=True,
         color=GREEN_DARK, letter_spacing=4)
add_text(p2 + 0.20, dd_top + 0.50, sub_w - 0.40, 0.38,
         "Coverage across 482 flowering-plant families",
         font=SERIF, size=14, italic=True, color=INK_SOFT)
add_image_contain(F("chart_families.png"),
                  p2 + 0.15, dd_top + HEAD_H,
                  sub_w - 0.30, dd_h - HEAD_H - 0.20)

# Panel 3: EDGE (right)
p3 = geo_left + 2 * (sub_w + GAP)
add_rect(p3, dd_top, sub_w, dd_h, CARD)
add_text(p3 + 0.20, dd_top + 0.12, sub_w - 0.40, 0.42,
         "MOST IRREPLACEABLE PLANTS", font=SANS, size=18, bold=True,
         color=GREEN_DARK, letter_spacing=4)
add_text(p3 + 0.20, dd_top + 0.50, sub_w - 0.40, 0.38,
         "Safeguarding the world’s most evolutionarily distinct flora",
         font=SERIF, size=14, italic=True, color=INK_SOFT)
add_image_contain(F("edge_irreplaceable.png"),
                  p3 + 0.15, dd_top + HEAD_H,
                  sub_w - 0.30, dd_h - HEAD_H - 0.20)

# ============================================================ PARTNERS + CTA
CTA_TOP = DD_BOT + 0.20
CTA_H = 2.40
add_rect(0, CTA_TOP, 48, CTA_H, GREEN_DARK)
add_rect(0, CTA_TOP, 48, 0.08, GOLD)

# Two logical zones: Partnerships (left half) + Institutional Home (right half).
# QR code moved to the header; no Join Us column needed here.
panel_h = CTA_H - 0.40
p_w = (48 - 2*0.9 - 0.6) / 2   # half of usable width
pa_l = 0.9
pb_l = pa_l + p_w + 0.6

# --- ZONE A: PARTNERS — WFO + EDGE logos closer together (top), then text
add_text(pa_l, CTA_TOP + 0.15, p_w, 0.40,
         "ACTIVE PARTNERSHIPS",
         font=SANS, size=14, bold=True, color=GOLD, letter_spacing=6)
part_card_y = CTA_TOP + 0.55
part_card_h = panel_h - 0.55
add_rect(pa_l, part_card_y, p_w, part_card_h, WHITE)
# Top row (logos, ~55% of card height) — both centred together, close gap
logo_row_h = part_card_h * 0.55
margin     = 0.08
logo_h     = logo_row_h - 2 * margin
wfo_aspect, edge_aspect = 1600 / 533, 291 / 142
wfo_w  = logo_h * wfo_aspect
edge_w = logo_h * edge_aspect
logo_gap = 0.30
total_lw = wfo_w + logo_gap + edge_w
logos_x  = pa_l + (p_w - total_lw) / 2
add_image_contain(F("wfo_logo_full_dark.png"),
                  logos_x, part_card_y + margin, wfo_w, logo_h)
add_image_contain(F("edge_logo_dark.png"),
                  logos_x + wfo_w + logo_gap, part_card_y + margin,
                  edge_w, logo_h)
# Bottom row: Anchored in the World's Herbaria text
text_y = part_card_y + logo_row_h + 0.05
text_h = part_card_h - logo_row_h - 0.10
add_text(pa_l + 0.20, text_y, p_w - 0.40, 0.40,
         "ANCHORED IN THE WORLD'S HERBARIA",
         font=SANS, size=12, bold=True, color=GREEN_DARK, letter_spacing=4)
add_paragraphs(
    pa_l + 0.20, text_y + 0.40, p_w - 0.40, text_h - 0.40,
    [
        ("", {"runs": [
            ("Our occurrences are anchored to ", {}),
            ("2,424 voucher specimens", {"bold": True}),
            (" deposited across ", {}),
            ("46 herbaria worldwide (and counting)", {"bold": True}),
            (" — top partners: Boise State University (1,771), Missouri "
             "Botanical Garden (253), MNHN Paris (56). We are actively "
             "digitising herbarium sheets to link images to our occurrences.",
             {}),
        ], "size": 12, "color": INK, "space_after": 0}),
    ],
    font=SERIF, line_spacing=1.22,
)

# --- ZONE B: INSTITUTIONAL HOME (BSU + Endowed Chair) ----------
add_text(pb_l, CTA_TOP + 0.15, p_w, 0.40,
         "INSTITUTIONAL HOME",
         font=SANS, size=14, bold=True, color=GOLD, letter_spacing=6)
b_card_h = part_card_h
bsu_logo_h = b_card_h - 0.20
bsu_logo_w = bsu_logo_h * (492 / 176)
bsu_card_x = pb_l
bsu_card_w = bsu_logo_w + 0.4
add_rect(bsu_card_x, part_card_y, bsu_card_w, b_card_h, WHITE)
slide.shapes.add_picture(F("bsu_logo_trim.png"),
                         Inches(bsu_card_x + 0.20),
                         Inches(part_card_y + (b_card_h - bsu_logo_h)/2),
                         Inches(bsu_logo_w), Inches(bsu_logo_h))
bt_l = bsu_card_x + bsu_card_w + 0.30
bt_w = p_w - (bt_l - pb_l)
add_text(bt_l, part_card_y + 0.00, bt_w, 0.55,
         "Dr. Christopher Davidson Endowed Chair in Botany",
         font=SERIF, size=18, bold=True, color=WHITE, line_spacing=1.10)
add_text(bt_l, part_card_y + 0.60, bt_w, b_card_h - 0.65,
         "Established in 2024 at Boise State University, with "
         "Associate Professor Sven Buerki as inaugural chair.",
         font=SERIF, size=13, italic=True, color=GOLD_SOFT, line_spacing=1.20)

# ============================================================ FOOTER
FT_TOP = CTA_TOP + CTA_H + 0.20
FT_H = 36 - FT_TOP
add_rect(0, FT_TOP, 48, FT_H, CREAM)
add_line(0, FT_TOP, 48, FT_TOP, color=GOLD, weight=2)

# References — compact font + tighter line spacing so all four lines fit
ref_l = 0.9
ref_w = 27.0
add_text(ref_l, FT_TOP + 0.10, ref_w, 0.35,
         "REFERENCES", font=SANS, size=13, bold=True,
         color=GREEN_DARK, letter_spacing=6)
add_line(ref_l, FT_TOP + 0.50, ref_l + 1.3, FT_TOP + 0.50, color=GOLD, weight=2)
ref_blocks = [
    ("", {"runs": [
        ("Forest, F., R. Brown, S. Buerki, et al. 2026. ", {}),
        ("“High risk of extinction across the flowering plant tree of life.”",
         {"italic": True}),
        (" Science 392 (6798): 655–659. doi.org/10.1126/science.adz0773", {}),
    ], "size": 11, "color": INK, "space_after": 2}),
    ("", {"runs": [
        ("Forest, F., J. Moat, E. Baloch, et al. 2018. ", {}),
        ("“Gymnosperms on the EDGE.”", {"italic": True}),
        (" Scientific Reports 8: 6053. doi.org/10.1038/s41598-018-24365-4", {}),
    ], "size": 11, "color": INK, "space_after": 2}),
    ("", {"runs": [
        ("Taylor, C. M., R. E. Gereau, W. D. Stevens, S. Buerki, "
         "O. M. Montiel, and S. Christoph. 2023. ", {}),
        ("“In Memoriam: Chris Davidson (1944–2022), the Idaho Botanist "
         "Who Botanized the World.”", {"italic": True}),
        (" Annals of the Missouri Botanical Garden 108. doi.org/10.3417/2023858", {}),
    ], "size": 11, "color": INK, "space_after": 2}),
    ("", {"runs": [
        ("Full publication list at ", {}),
        ("floraoftheworld.org/publications", {"bold": True, "color": GREEN_DARK}),
    ], "size": 11, "italic": True, "color": INK_SOFT, "space_after": 0}),
]
add_paragraphs(ref_l, FT_TOP + 0.55, ref_w, max(FT_H - 0.60, 0.10),
               ref_blocks, font=SERIF, line_spacing=1.12)

# --- DOWNLOAD THIS POSTER (centre of footer): label on top, QR + URL beside ---
dl_l = 28.5
dl_w = 5.8
add_text(dl_l, FT_TOP + 0.10, dl_w, 0.30,
         "DOWNLOAD THIS POSTER",
         font=SANS, size=12, bold=True, color=GREEN_DARK,
         letter_spacing=4, align=PP_ALIGN.CENTER)
# QR sits below the label on the left; URL text sits to its right
qr_size_dl = 1.12
qr_top_dl  = FT_TOP + 0.45
qr_l_dl    = dl_l + 0.10
slide.shapes.add_picture(F("qr_github.png"),
                         Inches(qr_l_dl), Inches(qr_top_dl),
                         Inches(qr_size_dl), Inches(qr_size_dl))
add_text(qr_l_dl + qr_size_dl + 0.20,
         qr_top_dl + qr_size_dl / 2 - 0.15,
         dl_w - qr_size_dl - 0.30, 0.30,
         "github.com/Flora-of-the-World/Botany_2026",
         font=SANS, size=10, italic=True, color=INK_SOFT,
         align=PP_ALIGN.LEFT)

# Contact (right) — compact for the 1.4" footer
ct_l = 35.5
ct_w = 47.0 - ct_l
add_text(ct_l, FT_TOP + 0.15, ct_w, 0.35,
         "CONTACT", font=SANS, size=13, bold=True,
         color=GREEN_DARK, letter_spacing=6, align=PP_ALIGN.RIGHT)
add_line(ct_l + ct_w - 1.1, FT_TOP + 0.55, ct_l + ct_w,
         FT_TOP + 0.55, color=GOLD, weight=2)
add_text(ct_l, FT_TOP + 0.65, ct_w, 0.45,
         "svenbuerki@boisestate.edu",
         font=SERIF, size=17, color=INK, align=PP_ALIGN.RIGHT)
add_text(ct_l, FT_TOP + 1.10, ct_w, max(FT_H - 1.20, 0.10),
         "Dr. Christopher Davidson Endowed Chair in Botany",
         font=SANS, size=12, italic=True, color=INK_SOFT,
         align=PP_ALIGN.RIGHT, line_spacing=1.25)

# ============================================================ SAVE
out_path = os.path.join(HERE, "Flora_of_the_World_Botany2026_Poster.pptx")
prs.save(out_path)
print("WROTE:", out_path)
